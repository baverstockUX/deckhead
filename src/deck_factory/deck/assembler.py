"""
Deck assembler for Deckhead.

Creates PowerPoint presentations using python-pptx with full-bleed images,
text overlays, and speaker notes.
"""

from pathlib import Path
from typing import List
import io

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

from ..core.models import DeckStructure, SlideContent, GeneratedImage, TextContent
from ..core.exceptions import AssemblyError, SlideCreationError, ImageProcessingError


class DeckAssembler:
    """
    Assembles PowerPoint presentations from structured content and images.

    Creates 16:9 slides with full-bleed images, text overlays, and speaker notes.
    """

    def __init__(self):
        """Initialize deck assembler with standard 16:9 dimensions."""
        # 16:9 dimensions in inches
        self.slide_width = Inches(13.333)
        self.slide_height = Inches(7.5)

    def create_deck(
        self,
        structure: DeckStructure,
        images: List[GeneratedImage],
        output_path: Path
    ) -> Path:
        """
        Create complete PowerPoint presentation.

        Args:
            structure: Deck structure with slide content
            images: Generated images for slides
            output_path: Path where to save the presentation

        Returns:
            Path to created presentation file

        Raises:
            AssemblyError: If deck creation fails
        """
        try:
            # Create presentation
            prs = Presentation()

            # Set slide dimensions to 16:9
            prs.slide_width = self.slide_width
            prs.slide_height = self.slide_height

            # Map images by slide number for quick lookup
            image_map = {img.slide_number: img for img in images}

            # Create slides
            for slide_content in structure.slides:
                # Get corresponding image
                image = image_map.get(slide_content.slide_number)
                if not image:
                    raise AssemblyError(
                        f"Missing image for slide {slide_content.slide_number}"
                    )

                # Create slide
                self._create_slide(prs, slide_content, image)

            # Ensure output directory exists
            output_path.parent.mkdir(parents=True, exist_ok=True)

            # Save presentation
            prs.save(str(output_path))

            return output_path

        except Exception as e:
            raise AssemblyError(f"Failed to create deck: {e}")

    def _create_slide(
        self,
        prs: Presentation,
        slide_content: SlideContent,
        image: GeneratedImage
    ) -> None:
        """
        Create a single slide with image and text overlay.

        Args:
            prs: Presentation object
            slide_content: Slide content definition
            image: Generated image for this slide

        Raises:
            SlideCreationError: If slide creation fails
        """
        try:
            # Add blank slide (layout index 6 is typically blank)
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # Route to appropriate layout handler based on layout_type
            layout_type = getattr(slide_content, 'layout_type', 'image-only')

            if layout_type == "split-left":
                self._create_split_left_slide(slide, slide_content, image.image_data)
            elif layout_type == "split-right":
                self._create_split_right_slide(slide, slide_content, image.image_data)
            elif layout_type == "panel":
                self._create_panel_slide(slide, slide_content, image.image_data)
            elif layout_type == "overlay":
                self._create_overlay_slide(slide, slide_content, image.image_data)
            else:  # Default: image-only
                self._add_full_bleed_image(slide, image.image_data)

            # Add speaker notes if present
            if slide_content.speaker_notes:
                self._add_speaker_notes(slide, slide_content.speaker_notes)

        except Exception as e:
            raise SlideCreationError(
                f"Failed to create slide {slide_content.slide_number}: {e}"
            )

    def _add_full_bleed_image(self, slide, image_data: bytes) -> None:
        """
        Add full-bleed image to slide.

        Args:
            slide: Slide object
            image_data: Raw image data bytes

        Raises:
            ImageProcessingError: If image insertion fails
        """
        try:
            # Convert bytes to file-like object
            image_stream = io.BytesIO(image_data)

            # Add picture at (0,0) with full slide dimensions
            slide.shapes.add_picture(
                image_stream,
                left=0,
                top=0,
                width=self.slide_width,
                height=self.slide_height
            )

        except Exception as e:
            raise ImageProcessingError(f"Failed to add image to slide: {e}")

    def _add_overlay_text(
        self,
        slide,
        text: str,
        position: str = "bottom"
    ) -> None:
        """
        Add text overlay with semi-transparent background.

        Args:
            slide: Slide object
            text: Text to overlay
            position: Position on slide ("bottom", "top", "center")
        """
        # Calculate position based on parameter
        if position == "bottom":
            left = Inches(0.5)
            top = Inches(6.0)
            width = Inches(12.333)
            height = Inches(1.0)
        elif position == "top":
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(12.333)
            height = Inches(1.0)
        else:  # center
            left = Inches(0.5)
            top = Inches(3.0)
            width = Inches(12.333)
            height = Inches(1.5)

        # Add text box
        textbox = slide.shapes.add_textbox(left, top, width, height)

        # Set semi-transparent background
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
        textbox.fill.transparency = 0.3  # 30% transparent

        # Configure text frame
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Add text
        p = text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER

        # Format text
        font = p.font
        font.name = 'Arial'
        font.size = Pt(36)
        font.bold = True
        font.color.rgb = RGBColor(255, 255, 255)  # White text

    def _add_speaker_notes(self, slide, notes: str) -> None:
        """
        Add speaker notes to slide.

        Args:
            slide: Slide object
            notes: Speaker notes text
        """
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame

        # Clear existing notes
        text_frame.clear()

        # Add notes
        text_frame.text = notes

    def _create_split_left_slide(
        self,
        slide,
        slide_content: SlideContent,
        image_data: bytes
    ) -> None:
        """
        Create slide with image on left half, text content on right half.

        Args:
            slide: Slide object
            slide_content: Slide content definition
            image_data: Raw image data bytes
        """
        # Image: Left half (6.667" x 7.5")
        image_stream = io.BytesIO(image_data)
        slide.shapes.add_picture(
            image_stream,
            left=0,
            top=0,
            width=Inches(6.667),
            height=Inches(7.5)
        )

        # Text content box: Right half
        text_left = Inches(6.667)
        text_top = Inches(0.75)
        text_width = Inches(6.0)
        text_height = Inches(6.0)

        textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)

        # White semi-transparent background for readability
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        textbox.fill.transparency = 0.05

        # Populate with structured content
        if slide_content.text_content:
            self._populate_text_content(textbox.text_frame, slide_content.text_content)

    def _create_split_right_slide(
        self,
        slide,
        slide_content: SlideContent,
        image_data: bytes
    ) -> None:
        """
        Create slide with text content on left half, image on right half.

        Args:
            slide: Slide object
            slide_content: Slide content definition
            image_data: Raw image data bytes
        """
        # Text content box: Left half
        text_left = Inches(0.667)
        text_top = Inches(0.75)
        text_width = Inches(6.0)
        text_height = Inches(6.0)

        textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)

        # White semi-transparent background for readability
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        textbox.fill.transparency = 0.05

        # Populate with structured content
        if slide_content.text_content:
            self._populate_text_content(textbox.text_frame, slide_content.text_content)

        # Image: Right half (6.667" x 7.5")
        image_stream = io.BytesIO(image_data)
        slide.shapes.add_picture(
            image_stream,
            left=Inches(6.667),
            top=0,
            width=Inches(6.667),
            height=Inches(7.5)
        )

    def _create_panel_slide(
        self,
        slide,
        slide_content: SlideContent,
        image_data: bytes
    ) -> None:
        """
        Create slide with full-width image on top, text panel below.

        Args:
            slide: Slide object
            slide_content: Slide content definition
            image_data: Raw image data bytes
        """
        # Image: Top 2/3 (13.333" x 5")
        image_stream = io.BytesIO(image_data)
        slide.shapes.add_picture(
            image_stream,
            left=0,
            top=0,
            width=Inches(13.333),
            height=Inches(5.0)
        )

        # Text panel: Bottom 1/3
        panel_left = Inches(0.5)
        panel_top = Inches(5.25)
        panel_width = Inches(12.333)
        panel_height = Inches(2.0)

        textbox = slide.shapes.add_textbox(panel_left, panel_top, panel_width, panel_height)

        # Light background for contrast
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray

        # Populate with structured content
        if slide_content.text_content:
            self._populate_text_content(textbox.text_frame, slide_content.text_content)

    def _create_overlay_slide(
        self,
        slide,
        slide_content: SlideContent,
        image_data: bytes
    ) -> None:
        """
        Create slide with minimal text overlaid on full-bleed image.

        Args:
            slide: Slide object
            slide_content: Slide content definition
            image_data: Raw image data bytes
        """
        # Full-bleed image
        self._add_full_bleed_image(slide, image_data)

        # Minimal overlay text (reuse existing _add_overlay_text method)
        if slide_content.text_content and slide_content.text_content.paragraphs:
            text = slide_content.text_content.paragraphs[0]
            self._add_overlay_text(slide, text, position="bottom")

    def _populate_text_content(self, text_frame, content: TextContent) -> None:
        """
        Populate text frame with structured content.

        Args:
            text_frame: PowerPoint text frame object
            content: TextContent model with bullets, statistics, paragraphs, callouts
        """
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.3)
        text_frame.margin_right = Inches(0.3)
        text_frame.margin_top = Inches(0.3)
        text_frame.margin_bottom = Inches(0.3)

        # Add bullets
        if content.bullets:
            for i, bullet in enumerate(content.bullets):
                p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
                p.text = bullet
                p.level = 0
                p.font.name = 'Arial'
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(50, 50, 50)  # Dark gray
                p.space_after = Pt(10)

        # Add statistics (formatted distinctly)
        if content.statistics:
            for stat in content.statistics:
                p = text_frame.add_paragraph()
                p.text = f"{stat['label']}: "
                p.font.name = 'Arial'
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(100, 100, 100)  # Medium gray

                # Add value in bold
                run = p.add_run()
                run.text = stat['value']
                run.font.bold = True
                run.font.size = Pt(24)
                run.font.color.rgb = RGBColor(0, 102, 204)  # Blue accent
                p.space_after = Pt(12)

        # Add paragraphs
        if content.paragraphs:
            for para_text in content.paragraphs:
                p = text_frame.add_paragraph()
                p.text = para_text
                p.font.name = 'Arial'
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(60, 60, 60)
                p.space_after = Pt(10)
                p.line_spacing = 1.2

        # Add callouts (with visual distinction)
        if content.callouts:
            for callout in content.callouts:
                # Title
                p = text_frame.add_paragraph()
                p.text = callout['title']
                p.font.name = 'Arial'
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(204, 102, 0)  # Orange accent
                p.space_after = Pt(4)

                # Text
                p = text_frame.add_paragraph()
                p.text = callout['text']
                p.font.name = 'Arial'
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(60, 60, 60)
                p.space_after = Pt(15)

    def create_title_slide(
        self,
        prs: Presentation,
        title: str,
        subtitle: str = None
    ) -> None:
        """
        Create a title slide (optional, for first slide).

        Args:
            prs: Presentation object
            title: Presentation title
            subtitle: Optional subtitle
        """
        # Use title slide layout (usually index 0)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title

        # Set subtitle if provided
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
